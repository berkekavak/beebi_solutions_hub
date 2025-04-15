import streamlit as st
import os
import torch
import pandas as pd
import json
import csv
from PIL import Image
from transformers import CLIPProcessor, CLIPModel
from ultralytics import YOLO
from pptx import Presentation
from streamlit_carousel import carousel
import base64

# Function to convert image to base64
def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode('utf-8')

# Convert the logo image to base64
logo_base64 = get_base64_image("images/BeeBI-Logo-Orj.png")

# Set page config to wide layout to utilize more width
st.set_page_config(layout="wide")

# Function for extracting images from PPTX
def extract_images_from_pptx(pptx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    prs = Presentation(pptx_path)
    img_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image") and shape.image:
                img = shape.image
                img_bytes = img.blob
                img_path = os.path.join(output_folder, f'image_{img_count}.png')
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                img_count += 1
    return output_folder

# CSS and header (unchanged)
# Updated CSS and header with centered logo and dynamic title
st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&display=swap');
body, html, * {{
    font-family: 'Poppins', Arial, sans-serif !important;
}}
.stApp, h1, h2, h3, h4, h5, h6, p, div, span, button, a {{
    font-family: 'Poppins', Arial, sans-serif !important;
}}
/* Fixed header for all pages */
.fixed-header {{
    position: fixed;
    top: 0;
    left: 0;
    width: 100%;
    background-color: #ffffff;
    z-index: 1000;
    padding: 10px 20px;
    display: flex;
    align-items: center;
    justify-content: center; /* Center the entire header content */
    box-shadow: 0 2px 10px rgba(0, 0, 0, 0.1);
    height: 60px;
}}
.fixed-header .logo-container {{
    position: absolute; /* Center the logo absolutely within the header */
    left: 50%;
    transform: translateX(-50%);
    z-index: 1001;
}}
.fixed-header .logo-container img {{
    height: 50px; /* Adjust the height of the logo image */
    vertical-align: middle;
}}
.fixed-header .header-title-container {{
    position: absolute; /* Position the title below the logo */
    top: 50px; /* Adjust based on header height and logo size */
    left: 50%;
    transform: translateX(-50%);
    width: 100%;
    text-align: center;
    z-index: 1000;
}}
.fixed-header .header-title {{
    font-family: 'Poppins', Arial, sans-serif !important;
    font-size: 1.5em;
    font-weight: 600;
    color: #000000;
    margin: 0;
    text-align: center;
}}
/* Adjust the app padding to account for the header */
.stApp {{
    padding-top: 100px; /* Increased padding to account for header and title */
    background-color: #f7f9fe;
}}
/* Center the main page title (Streamlit's default title) */
div[data-testid="stMarkdownContainer"] h1:not(.solution-title),
div[data-testid="stMarkdownContainer"] h2:not(.solution-title) {{
    font-family: 'Poppins', Arial, sans-serif !important;
    font-size: 3.5em;
    font-weight: 500;
    text-align: center !important;
    width: 100%;
}}
/* Ensure all markdown text (body text) uses Poppins and has readable color */
div[data-testid="stMarkdownContainer"] p,
div[data-testid="stMarkdownContainer"] {{
    font-family: 'Poppins', Arial, sans-serif !important;
    color: #333333;
}}
/* Ensure solution page body text uses Poppins and has readable color */
.solution-content p {{
    font-family: 'Poppins', Arial, sans-serif !important;
    font-size: 1em;
    color: #333333;
    line-height: 1.6;
    margin: 20 0 0 0px;
    width: 130%;
}}
.solution-content h3 {{
    font-size: 1.2em;
    font-weight: 600;
    color: #1a2942;
    margin: 0 0 10px 0;
}}
#MainMenu {{visibility: hidden;}}
footer {{visibility: hidden;}}
header {{visibility: hidden;}}
h1.solution-title {{
    font-family: 'Poppins', Arial, sans-serif !important;
    font-size: 3em;
    font-weight: 500;
    color: #000000;
    margin-bottom: 20px;
    text-align: center !important;
    line-height: 1.2;
    letter-spacing: -0.02em;
    display: block;
    width: 100%;
}}
.solution-container {{
    display: flex;
    align-items: center;
    margin: 80px;
    width: 100%;
    max-width: 1024px;
}}
.solution-text {{
    width: 70%; 
    padding-right: 0px;
}}
.solution-point {{
    display: flex;
    align-items: flex-start;
    margin-bottom: 30px;
}}
.solution-number {{
    font-size: 2em;
    font-weight: 700;
    color: #d3e1f5;
    margin-right: 20px;
}}
.solution-image {{
    width: 150%; 
    border-radius: 10px;
    box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    padding-right: 200px;
}}
.back-button {{
    background-color: #FADB49;
    color: #1a2942;
    border: none;
    padding: 10px 20px;
    border-radius: 5px;
    font-weight: 600;
    cursor: pointer;
    transition: all 0.3s ease;
}}
.back-button:hover {{
    background-color: #e6c73e;
}}
.hexagon {{
    width: 100px;
    height: 100px;
    background: linear-gradient(145deg, #ffffff 0%, #f7f9fe 100%);
    clip-path: polygon(50% 0%, 100% 25%, 100% 75%, 50% 100%, 0% 75%, 0% 25%);
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    padding: 0px;
    box-sizing: border-box;
    text-align: center;
    box-shadow: 0 4px 15px rgba(0,0,0,0.1), 0 8px 20px rgba(0,0,0,0.05), inset 0 -2px 5px rgba(0,0,0,0.05);
    border: 10px solid rgba(255,255,255,0.8);
    transition: all 0.3s ease;
    cursor: pointer;
}}
.hexagon:hover {{
    background: linear-gradient(145deg, #1a2942 0%, #2c4268 100%);
    transform: scale(1.1);
    box-shadow: 0 8px 25px rgba(0,0,0,0.15), 0 12px 30px rgba(0,0,0,0.1), inset 0 -3px 8px rgba(0,0,0,0.15);
}}
.hexagon h3 {{
    margin: 0;
    font-size: 1em;
    color: #1a2942;
}}
.hexagon:hover h3 {{
    color: white;
}}
.hexagon-container {{
    display: flex;
    flex-direction: column;
    align-items: center;
    margin: 0 auto;
    width: 100%;
    max-width: 1024px;
    padding: 0px !important;
}}
.hexagon-row {{
    display: flex;
    justify-content: center;
    align-items: center;
    margin-bottom: -80px !important;
}}
div[data-testid="stHorizontalBlock"] {{
    gap: 0px !important;
    display: flex;
    justify-content: center;
}}
.stButton>button {{
    width: 200px;
    height: 200px;
    background: linear-gradient(145deg, #ffffff 0%, #f7f9fe 100%);
    clip-path: polygon(50% 0%, 100% 25%, 100% 75%, 50% 100%, 0% 75%, 0% 25%);
    padding: 0px;
    margin: -10px -10px !important;
    box-shadow: 0 4px 15px rgba(0,0,0,0.1), 0 8px 20px rgba(0,0,0,0.05), inset 0 -2px 5px rgba(0,0,0,0.05);
    border: 0px solid rgba(0,0,0,0.1);
    transition: all 0.3s ease;
    cursor: pointer;
    display: flex;
    align-items: center;
    justify-content: center;
    text-align: center;
    filter: drop-shadow (3px 3px 3px rgba(1 1 1 1));
}}
.stButton>button:hover {{
    background: #FADB49;
    transform: scale(1.1);
    box-shadow: 0 8px 25px rgba(0,0,0,0.15), 0 12px 30px rgba(0,0,0,0.1), inset 0 -3px 8px rgba(0,0,0,0.15);
}}
.stButton>button h3 {{
    margin: 0;
    font-size: 1em;
    color: #1a2942;
}}
.stButton>button:hover h3 {{
    color: black;
}}
div[data-testid="column"] {{
    padding: 0px !important;
    margin: 0px !important;
    display: flex;
    justify-content: center;
    align-items: center;
}}
.row-offset {{
    position: relative;
    left: 55px; /* Offset odd rows to create honeycomb pattern */
}}
.row-container {{    
    margin-bottom: -60px !important; /* Pull rows closer together */
}}
.hexagon-text {{ 
    max-width: 90%;
    text-align: center;
}}
</style>

<div class="fixed-header">
    <div class="logo-container">
        <img src="data:image/png;base64,{logo_base64}" alt="BeeBI Logo">
    </div>
</div>
""", unsafe_allow_html=True)



# Main page (unchanged)
def main_page():
    st.markdown("""
        <p style="font-size: 2.5em; color: #333; text-align: center;">
            <strong>Solutions Hub</strong>
        </p>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style="max-width: 1100px; margin: 1px 1px; padding: 0 -40px; text-align: center;">
    """, unsafe_allow_html=True)

    st.markdown("""
        <p style="font-size: 1.2em; color: #333; text-align: center;">
            Discover BeeBI’s innovative <strong>AI</strong> and <strong>data solutions</strong> through interactive demos and detailed descriptions, showcasing the power of our <strong>technology</strong>. Learn more about our offerings at <a href="https://www.beebi-consulting.com/" target="_blank">BeeBI Consulting</a>.
        </p>
    """, unsafe_allow_html=True)

    solutions = [
        {"id": "assortment-analytics", "title": "Assortment Analytics", "link": "assortment_analytics"},
        {"id": "demand-prediction", "title": "Demand Prediction", "link": "demand_prediction"},
        {"id": "price-elasticity", "title": "Price Elasticity", "link": "price_elasticity"},
        {"id": "material-product-lifecycle", "title": "Material/Product Lifecycle", "link": "material_product_lifecycle"},
        {"id": "inventory-flow-optimization", "title": "Inventory Flow Efficiency", "link": "inventory_flow_optimization"},
        {"id": "competitive-intelligence", "title": "Competitive Intelligence", "link": "competitive_intelligence"},
        {"id": "markdown-optimization", "title": "Markdown Optimization", "link": "markdown_optimization"},
        {"id": "data-sourcing", "title": "Data Readiness Status", "link": "data_sourcing"},
        {"id": "image-attribute-prediction", "title": "Image Extraction & Prediction", "link": "image_attribute_prediction"},
        {"id": "football-player-recommender", "title": "Player Recommender", "link": "football_player_recommender"},
        {"id": "topic-summarization", "title": "Topic Summarization & Grouping", "link": "topic_summarization"},
        {"id": "product-article-analyzer", "title": "Product/Article Analyzer", "link": "product_article_analyzer"},
    ]

    with st.spinner("Loading Use Cases..."):
        # Row 1: 3 hexagons
        row1 = st.columns([1, 1, 1, 1, 1])
        for i in range(3):
            with row1[i + 1]:
                solution = solutions[i]
                if st.button(f"### {solution['title']}", key=solution['id']):
                    st.session_state.page = solution['link']
                    st.rerun()

        # Row 2: 4 hexagons
        row2 = st.columns([1, 1, 1, 1])
        for i in range(4):
            with row2[i]:
                solution = solutions[i + 3]
                if st.button(f"### {solution['title']}", key=solution['id']):
                    st.session_state.page = solution['link']
                    st.rerun()

        # Row 3: 3 hexagons
        row3 = st.columns([1, 1, 1, 1, 1])
        for i in range(3):
            with row3[i + 1]:
                solution = solutions[i + 7]
                if st.button(f"### {solution['title']}", key=solution['id']):
                    st.session_state.page = solution['link']
                    st.rerun()

        # Row 4: 2 hexagons
        row4 = st.columns([1, 1, 1, 1])
        for i in range(2):
            with row4[i + 1]:
                solution = solutions[i + 10]
                if st.button(f"### {solution['title']}", key=solution['id']):
                    st.session_state.page = solution['link']
                    st.rerun()

    st.markdown("</div>", unsafe_allow_html=True)

# Assortment Analytics page
def assortment_analytics_page():
    st.markdown('<h1 class="solution-title">Assortment Analytics</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>Managing the complexity of high-volume assortment planning processes within a corporate framework demands intense consideration of various dimensions and key performance indicators. The establishment and maintenance of an optimal assortment analytic platform constitute a formidable challenge for enterprises. This complexity underscores the pivotal role such a system plays in contributing significantly to a company’s overarching success and strategic objectives, making this a critical problem to be solved for organizations across industries.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>A robust Single Source of Truth Platform integrates ranging progresses, sales, forecasts, and key performance indicators, providing a centralized repository for high-volume assortment planning. It serves as the authoritative reference for cross-departmental decisions, contributing significantly to a company’s success. The platform is designed for continuous improvement to adapt to evolving needs.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>The platform streamlines cross-departmental coordination, enhancing speed and efficiency in assortment planning processes. It ensures a trusted reference for decisions, contributing significantly to overall company success. The continuously improved platform adapts to evolving needs and supports a better decision-making mechanism that increases value across products in the market.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        # Path to the assortment-analytics folder
        solution_folder = "images/assortment-analytics"
        # Get list of image files in the folder
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        # Create carousel items dynamically
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        # Display carousel if there are images, otherwise show a placeholder
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/assortment-analytics folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Competitive Intelligence page
def competitive_intelligence_page():
    st.markdown('<h1 class="solution-title">Competitive Intelligence</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The business challenge is the need for competitive intelligence to analyze the market comprehensively. The problem involves leveraging web scraping data to enable monitoring of current and historical market analyses, including competitive brands and products. The importance lies in the ability to accurately define and strategically position products and brands based on real-time and historical market insights. This entails staying informed about market dynamics, competitor strategies, and product performance, which is crucial for making informed decisions, staying competitive, and maximizing market share. A solution is vital for navigating the ever-changing market landscape and ensuring a strong market position and pricing strategy.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution to competitive intelligence involves implementing a robust web scraping system with an integrated alerting mechanism. This system enables comprehensive monitoring of current and historical market analyses, providing insights into competitive brands and products. The alerting mechanism ensures real-time notifications for significant market changes, allowing businesses to promptly adjust strategies. This comprehensive approach empowers businesses to accurately define and strategically position their products and brands based on dynamic market insights, enhancing competitiveness and maximizing market share.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution, combining advanced web scraping with a proactive alerting mechanism, delivers unparalleled value in competitive intelligence. By providing comprehensive insights into market analyses, including competitive brands and products, it empowers informed decision-making. The alerting mechanism ensures timely notifications of market changes, enabling prompt strategic adjustments. This enhances accuracy in  in positioning products and brands, supports competitiveness, and maximizes market share, driving sustained business success through strategic agility and market responsiveness.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        # Path to the competitive-intelligence folder
        solution_folder = "images/competitive-intelligence"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/competitive-intelligence folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Data Sourcing page
def data_sourcing_page():
    st.markdown('<h1 class="solution-title">Data Sourcing</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>Managing product attributes at scale presented a significant challenge for a retail client with extensive product catalogs. Before each new season, thousands of products required complete specifications including color, length, material, and other critical attributes. Without these specifications, production and distribution couldn’t proceed. The manual attribute entry process lacked visibility into completion status, creating bottlenecks in product readiness. Teams had no systematic way to identify missing critical attributes or ensure entered data met quality standards, resulting in delayed launches and operational inefficiencies.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>We developed a comprehensive analytics platform that provides visibility into attribute completion status across the entire product catalog. The solution includes: <br>- Dynamic dashboards showing attribute completion percentages by product category<br>- Prioritized lists of products missing critical attributes<br>- Automated data quality validation that flags incorrect attribute formats (e.g., numeric values in color fields)<br>- Team-specific views allowing departments to focus on their relevant attributes<br>- Progress tracking that measures completion rates against seasonal timelines</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>The Attribute Management Analytics Platform transformed the client’s product preparation process, delivering substantial benefits: <br>- Reduced time-to-market by 40% through streamlined attribute completion<br>- Improved resource allocation by focusing teams on critical missing attributes<br>- Eliminated quality issues from incorrectly formatted attributes, reducing rework and production delays<br>- Enhanced cross-departmental coordination with shared visibility into attribute completion status<br>- Enabled data-driven decision making about product prioritization and seasonal readiness<br>- Increased operational efficiency by replacing manual attribute checking with automated systems</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/data-sourcing"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/data-sourcing folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Demand Prediction page
def demand_prediction_page():
    st.markdown('<h1 class="solution-title">Demand Prediction</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The business problem of demand prediction centers on the challenge of accurately forecasting customer purchases over a specific timeframe. This is critical for businesses to optimize inventory management, minimize stock-related issues, and enhance overall supply chain efficiency. The problem involves leveraging historical data, market trends, and advanced analytics to make precise predictions, ensuring that businesses can align resources with anticipated customer demand, reduce costs, and maintain effective operations.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>To address the challenge of accurate demand prediction, our solution integrates advanced analytics and machine learning techniques. The solution empowers organizations to align resources with anticipated customer demand, resulting in reduced costs and more effective operations. With a focus on data-driven insights, our platform ensures businesses can make informed decisions to meet customer expectations and stay ahead in the dynamic market landscape.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution for accurate demand prediction delivers tangible value to businesses. By integrating advanced analytics and machine learning, it provides precise forecasting, enabling optimized inventory management and enhanced supply chain efficiency. This reduces costs, improves operational effectiveness, and ensures businesses can meet customer demands in a dynamic market landscape.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/demand-prediction"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/demand-prediction folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Football Player Recommender page
def football_player_recommender_page():
    st.markdown('<h1 class="solution-title">Football Player Recommender</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The challenge lies in the inefficient process of scouting and recommending football players, which often lacks data-driven insights and alignment with team strategies. Traditional scouting methods rely heavily on manual analysis, leading to inconsistent player evaluations, therefore undesired transfers. There’s a critical need for a solution that leverages data to provide objective recommendations, ensuring players align with a team’s playstyle and strategic needs while optimizing the scouting process.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution, the Football Player Recommender, addresses the inefficiencies in scouting by leveraging data-driven insights to recommend players based on performance metrics, playstyle, and tactical compatibility. We propose a system that analyzes player data to provide tailored recommendations, ensuring alignment with a team’s strategy. This solution streamlines the scouting process, offering an objective and efficient approach to player selection. By using advanced data analytics, the Football Player Recommender enhances the accuracy and relevance of player recommendations for teams.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>The Football Player Recommender delivers significant value by providing data-driven recommendations, fostering effective and strategic scouting processes. The system’s ability to analyze performance metrics and ensure tactical compatibility offers an objective approach, improving recruitment efficiency. This reduces the risk of poor player selections, optimizes team performance, and saves time in the scouting process. The value lies in the enhanced decision-making, cost-effective recruitment, and improved team alignment, ultimately contributing to a stronger competitive edge on the field.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/football-player-recommender"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/football-player-recommender folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

    st.write("Leverage data-driven insights to scout and recommend football players based on performance metrics, playstyle, and tactical compatibility.")

# Image Attribute Prediction page
def image_attribute_prediction_page():
    st.markdown('<h1 class="solution-title">Image Attribute Prediction</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The challenge lies in the manual process of extracting and classifying product attributes from PPTX files, which prevents efficient product categorization in retail and e-commerce. Businesses rely on labor-intensive methods to extract images, apply object detection, and extract attributes one by one, leading to delays in product classification, operational inefficiencies, and reduced searchability. There’s a need for a solution that automates this process, enabling faster and accurate attribute prediction to streamline operations and improve efficiency in product management.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Image Attribute Prediction addresses the challenge of inefficient product attribute extraction by leveraging AI vision models to automate the process. We propose a system that extracts images from PPTX files and uses advanced AI vision models to predict product attributes, enabling faster classification and automation. This solution enhances the product categorization process, providing an efficient approach for retail and e-commerce businesses. By automating attribute prediction, our solution eliminates bottlenecks and enhances the accuracy and speed of product management workflows.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution delivers significant value by automating the classification of product attributes, fostering efficient product management processes. The use of AI vision models to predict attributes from PPTX files provides a streamlined approach, enhancing operational efficiency in retail and e-commerce. This reduces manual effort and operational costs while improving product searchability, ensuring businesses can quickly categorize products and meet customer demands. The value lies in the improved efficiency, cost reduction, and enhanced categorization capabilities, contributing to better customer experiences and a stronger competitive edge.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/image-attribute-prediction"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/image-attribute-prediction folder.")

    if st.button("Back to BeeBI Solution Hub", key="back_to_hub_image_attribute"):
        st.session_state.page = "main"
        st.rerun()

    uploaded_pptx = st.file_uploader("Upload a PPTX file", type=["pptx"])
    if uploaded_pptx:
        with st.spinner("Processing PPTX file..."):
            pptx_path = os.path.join("temp.pptx")
            with open(pptx_path, "wb") as f:
                f.write(uploaded_pptx.read())
            
            output_folder = "images_output"
            extract_images_from_pptx(pptx_path, output_folder)
            
            st.write("Extracted Images:")
            image_paths = [os.path.join(output_folder, img_name) for img_name in os.listdir(output_folder)]
            
            index = st.session_state.get("image_index", 0)
            
            if image_paths:
                image = Image.open(image_paths[index])
                st.image(image, caption=os.path.basename(image_paths[index]))
                
                with st.spinner("Predicting attributes..."):
                    model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
                    processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                    
                    attributes = {
                        "Product Type": ["T-shirt", "Hoodie", "Sweatshirt", "Jacket", "Pants", "Shorts", "Shoes", "Hat", "Bag"],
                        "Illustration Type": ["Product Drawing or Sketch", "Product Photo"],
                        "Color": ["Red", "Blue", "Black", "White", "Green", "Pink", "Yellow", "Purple", "Orange", "Brown", "Gray", "Beige", "Claret Red", "Multicolor"],
                        "Collar Type": ["Round neck", "V-neck", "Polo collar", "Turtleneck", "Buttoned collar"],
                        "Material": ["Cotton", "Polyester", "Blended fabric"],
                        "Sleeve Length": ["Short sleeves", "Long sleeves", "Sleeveless"],
                        "Lining": ["Unlined", "Partially lined", "Fully lined"],
                        "Pocket Type": ["Zipper", "Flap", "Patch", "Slit", "None"],
                        "Shoe Material": ["Leather", "Mesh", "Knit", "Suede", "Canvas", "Rubber", "Ethylene Vinyl Acetate", "Recycled Materials"],
                        "Logo Type": ["adidas Originals (Trefoil) Logo", "adidas Performance (Three Stripes) Logo", "No Logo"]
                    }
                    
                    def get_best_match(category, options, image):
                        inputs = processor(text=options, images=image, return_tensors="pt", padding=True)
                        with torch.no_grad():
                            outputs = model(**inputs)
                            logits_per_image = outputs.logits_per_image
                            probs = logits_per_image.softmax(dim=1).squeeze()
                        probs = probs if isinstance(probs, torch.Tensor) else torch.tensor([probs])
                        best_match_idx = probs.argmax().item()
                        return options[best_match_idx], dict(zip(options, probs.tolist()))
                    
                    st.header("Image Attributes")
                    probabilities = {}
                    for category in attributes.keys():
                        value, probs = get_best_match(category, attributes[category], image)
                        st.write(f"**{category}:** {value}")
                        probabilities[category] = probs
                    
                    with st.expander("Show All Probabilities"):
                        for category, probs in probabilities.items():
                            st.write(f"**{category} Probabilities:** {json.dumps(probs, indent=4)}")
                
                col1, col2 = st.columns([1, 1])
                with col1:
                    if st.button("Previous Image"):
                        st.session_state.image_index = (index - 1) % len(image_paths)
                        st.rerun()
                with col2:
                    if st.button("Next Image"):
                        st.session_state.image_index = (index + 1) % len(image_paths)
                        st.rerun()

# Inventory Flow Efficiency Optimization page
def inventory_flow_optimization_page():
    st.markdown('<h1 class="solution-title">Inventory Flow Efficiency Optimization</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The business faces a challenge in optimizing inventory flow for increased efficiency. The problem is the absence of a tailored algorithm to monitor and enhance production, stock, and logistics optimization on a per-product basis. This results in inefficiencies such as excessive stock levels, high logistics costs, and suboptimal product efficiency and profitability. Additionally, there is a need to address product waste. A sophisticated algorithm that considers operational parameters and exceptions is crucial for achieving precise and data-driven inventory flow management, minimizing waste, and optimizing costs for increased profitability.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution involves the development of a sophisticated algorithm designed for per-product inventory flow optimization. This tailored approach addresses inefficiencies related to excessive stock, high logistics costs, and suboptimal product efficiency. By considering operational parameters and exceptions for each product, our solution provides a precise and data-driven approach to inventory flow management. This streamlines processes, enabling businesses to operate with minimal excess stock, efficient logistics, and optimized production, significantly enhancing overall profitability by minimizing waste and reducing costs.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution introduces a sophisticated algorithm for per-product inventory flow optimization, addressing inefficiencies such as excess stock, high logistics costs, and suboptimal product efficiency. This provides precise and data-driven inventory management, streamlining processes and enabling efficient operations with minimal excess stock. The result is increased profitability through reduced waste and lower costs, offering a strategic approach to enhance operational efficiency and financial outcomes.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/inventory-flow-optimization"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/inventory-flow-optimization folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Markdown Optimization page
def markdown_optimization_page():
    st.markdown('<h1 class="solution-title">Markdown Optimization</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1])
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The challenge is optimizing markdown rates for each product during campaigns by automating the process based on defined parameters, exceptions, and input values. The problem is establishing an efficient system to automatically calculate and apply the best markdown rates, considering individual product characteristics and campaign dates. Solving this is crucial for maximizing campaign effectiveness, minimizing excess inventory, and strategically aligning pricing for enhanced profitability. A robust markdown optimization process streamlines decision-making, ensuring a systematic approach to pricing adjustments during campaigns.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our proposed solution to markdown optimization entails the development of a sophisticated algorithmic system. Through the integration of machine learning and data analytics, our platform automates the computation of optimal markdown rates for each product during campaign periods. This system considers defined parameters, exceptions, and input values to ensure a dynamic and individualized approach. Incorporating campaign start and end dates, along with other pertinent factors, our solution establishes a comprehensive process for determining the most effective markdown rates. This streamlines decision-making and enhances the efficacy of promotional campaigns by minimizing excess inventory and maximizing profitability.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution automates precise markdown rate calculations for each product during campaigns, leveraging advanced algorithms and data analytics. This ensures dynamic pricing adjustments, streamlining decision-making, reducing excess inventory, and maximizing profitability. The value lies in improved campaign effectiveness, cost efficiency, and enhanced profitability through a systematic pricing approach.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/markdown-optimization"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/markdown-optimization folder.")

        # Add the video below the carousel
        video_path = "MDO-Demo-20250302_171805-MeetingRecording.mp4"
        try:
            st.video(video_path)
        except Exception as e:
            st.error(f"Error loading video: {e}")
            st.write("Please ensure the video file 'MDO-Demo-20250302_171805-MeetingRecording.mp4' is in the correct directory.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Material/Product Lifecycle Management page
def material_product_lifecycle_page():
    st.markdown('<h1 class="solution-title">Material/Product Lifecycle Management</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>Effectively managing material or product life cycles involves implementing a comprehensive system for monitoring, strategic decision-making, and optimizing profitability for each product in the market. The challenge is to establish an efficient management system, crucial for aligning resources, minimizing waste, and making informed decisions to maximize profitability throughout the product life cycle.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution to material or product life cycle management involves the implementation of an integrated system. Utilizing advanced analytics and automation, our platform facilitates real-time monitoring, strategic decision-making, and profitability optimization for each product in the market. By providing comprehensive insights and analytics at every life cycle stage, our solution empowers businesses to make informed decisions, minimize waste, and align resources effectively. This systematic approach ensures the best possible life cycle management, maximizing profitability for each product.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution offers substantial value by providing an integrated system for life cycle management. Through advanced analytics and automation, the platform delivers real-time insights, facilitating informed decision-making and profitability optimization. The value lies in the ability to minimize waste, align resources effectively, and maximize profitability throughout the entire life cycle. This systematic approach enhances operational efficiency, reduces costs, and achieves the best possible outcomes for each product in the portfolio.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/material-product-lifecycle"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/material-product-lifecycle folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Price Elasticity page
def price_elasticity_page():
    st.markdown('<h1 class="solution-title">Price Elasticity</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The challenge involves determining price elasticity, a critical factor in establishing optimal pricing strategies for products during campaigns. Accurately assessing consumer sensitivity to price changes is essential for defining the most profitable pricing structure for each product, particularly during promotional periods. This problem’s significance lies in aligning pricing strategies with consumer behavior to maximize profitability and competitiveness. Precise understanding of price elasticity is imperative for setting prices that resonate with consumers, enhancing campaign effectiveness and achieving optimal financial outcomes.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution to the challenge of price elasticity involves implementing advanced analytics and machine learning algorithms. By leveraging historical pricing and sales data, our platform accurately assesses consumer sensitivity to price changes for each product. This data-driven approach enables the establishment of optimal pricing strategies during campaign periods, ensuring the highest profitability while maintaining competitiveness in the market. The solution empowers businesses to fine-tune pricing structures based on precise insights into consumer behavior, leading to effective and profitable promotional campaigns.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Our solution delivers substantial value by employing advanced analytics and machine learning to assess price elasticity for each product. This data-driven approach empowers businesses to optimize pricing strategies during campaigns, maximizing profitability and competitiveness. The precise insights into consumer behavior enable strategic adjustments in pricing structures, ensuring effective and profitable promotional campaigns. The solution facilitates informed decision-making, enhancing overall financial outcomes and market positioning for businesses.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/price-elasticity"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/price-elasticity folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Product Article Analyzer page
def product_article_analyzer_page():
    st.markdown('<h1 class="solution-title">Range Viewer: Product Article Analyzer</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The current reporting system focuses on quantitative data visualization, such as dashboards and tables, but end-users require a more intuitive and visual way to analyze their product range. They need images integrated into detailed grids to enhance understanding of product distribution and characteristics. The existing Range Viewer is no longer available, creating a gap in analytical capabilities. Based on power user feedback, the new solution must address key use cases, such as analyzing product distribution, visualizing product spread, retrieving relevant product details, and enabling head-to-head product comparisons.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>To meet these requirements, a new reporting solution will integrate images into the detailed grids of Range Insights. This report will provide four key functionalities: <br>- Analyzing product distribution across qualitative and quantitative dimensions, including seasonality<br>- Visualizing product spread based on selected attributes<br>- Offering quick access to relevant product information using images and attributes<br>- Enabling direct product comparisons with enriched data visualization<br>The report leverages pre-built text and image similarity scores (with AI) to streamline article similarity analysis, enhancing efficiency by eliminating custom development. Users will benefit from flexible data export options (CSV, Excel, PowerPoint) and quick access to the Article Similarity form through a notification button.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>This solution significantly improves user experience by providing a more intuitive and visually enriched way to analyze product data. By incorporating images into the reporting framework, users gain deeper insights into their product range, leading to better decision-making. The use of pre-validated similarity metrics accelerates workflows and enhances accuracy, ensuring stakeholder confidence. The ability to export data in multiple formats and access key tools seamlessly improves efficiency and streamlines operations across business units.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/product-article-analyzer"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/product-article-analyzer folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

# Topic Summarization and Grouping page
def topic_summarization_page():
    st.markdown('<h1 class="solution-title">Topic Summarization and Grouping (Generative AI)</h1>', unsafe_allow_html=True)

    col1, col2 = st.columns([1.2, 1], vertical_alignment="center")
    
    with col1:
        st.markdown("""
        <div class="solution-text">
            <div class="solution-point">
                <span class="solution-number">01</span>
                <div class="solution-content">
                    <h3>PROBLEM</h3>
                    <p>The challenge is the inefficient process of analyzing and organizing large volumes of text, which hinders productivity and insight extraction. Businesses and teams often struggle with unorganized data, leading to difficulties in identifying key topics, grouping related content, and deriving actionable insights. There’s a pressing need for a solution that can efficiently summarize and structure text data, enabling faster analysis and better decision-making while reducing manual effort.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">02</span>
                <div class="solution-content">
                    <h3>SOLUTION</h3>
                    <p>Our solution, Topic Summarization and Grouping (Generative AI), tackles the inefficiency of text analysis by using Generative AI to summarize key topics and group related content from large volumes of text. We propose a system that automatically analyzes and organizes unorganized data, providing concise summaries and structured groupings. This solution streamlines the process of extracting insights, offering an efficient approach to text management. By leveraging Generative AI, our system enhances the speed and accuracy of content analysis for businesses and teams.</p>
                </div>
            </div>
            <div class="solution-point">
                <span class="solution-number">03</span>
                <div class="solution-content">
                    <h3>VALUE</h3>
                    <p>Topic Summarization and Grouping (Generative AI) offers significant value by automating the analysis of large text datasets, boosting productivity and insight generation. The AI-driven summarization and grouping provide an accurate approach, enhancing overall efficiency in data management. This reduces manual effort, saves time, and enables teams to quickly extract actionable insights from unorganized data. The value lies in the improved productivity, reduced operational costs, and enhanced decision-making capabilities, empowering organizations to make faster, data-driven decisions.</p>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        solution_folder = "images/topic-summarization"
        image_files = [f for f in os.listdir(solution_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        
        test_items = [
            dict(
                title="",
                text="",
                img=os.path.join(solution_folder, img_file),
                link=""
            ) for img_file in image_files
        ]
        
        if test_items:
            carousel(items=test_items)
        else:
            st.write("No images found in images/topic-summarization folder.")

    if st.button("Back to BeeBI Solution Hub"):
        st.session_state.page = "main"
        st.rerun()

    st.write("Summarize key topics and group related content using Generative AI to analyze and organize large volumes of text efficiently.")

    text_input = st.text_area("Enter text to summarize and group:", height=200)
    
    if st.button("Summarize and Group"):
        with st.spinner("Processing text..."):
            if text_input:
                st.write("**Summary:**")
                st.write("This is a mock summary of the input text.")
                st.write("**Grouped Topics:**")
                st.write("- Topic 1: General theme of the text")
                st.write("- Topic 2: Secondary theme identified")
            else:
                st.write("Please enter some text to process.")

# Main app logic to handle navigation
if "page" not in st.session_state:
    st.session_state.page = "main"

if st.session_state.page == "main":
    main_page()
elif st.session_state.page == "image_attribute_prediction":
    image_attribute_prediction_page()
elif st.session_state.page == "football_player_recommender":
    football_player_recommender_page()
elif st.session_state.page == "topic_summarization":
    topic_summarization_page()
elif st.session_state.page == "demand_prediction":
    demand_prediction_page()
elif st.session_state.page == "assortment_analytics":
    assortment_analytics_page()
elif st.session_state.page == "product_article_analyzer":
    product_article_analyzer_page()
elif st.session_state.page == "markdown_optimization":
    markdown_optimization_page()
elif st.session_state.page == "data_sourcing":
    data_sourcing_page()
elif st.session_state.page == "material_product_lifecycle":
    material_product_lifecycle_page()
elif st.session_state.page == "price_elasticity":
    price_elasticity_page()
elif st.session_state.page == "inventory_flow_optimization":
    inventory_flow_optimization_page()
elif st.session_state.page == "competitive_intelligence":
    competitive_intelligence_page()

# Add the footer
st.write("---")
st.write("Copyright © 2025 BeeBI Consulting GmbH. All rights reserved")